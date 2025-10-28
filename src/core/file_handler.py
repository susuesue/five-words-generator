"""
文件处理模块
负责从各种文档格式中提取文本内容
"""

from pathlib import Path

import PyPDF2
from docx import Document
import pandas as pd

from src.utils.config import (
    SUPPORTED_FORMATS,
    MAX_PDF_PAGES,
    MAX_DOCX_PARAGRAPHS,
    MAX_XLSX_ROWS,
    MAX_TXT_CHARS,
    MAX_PPTX_SLIDES
)
from src.utils.logger import get_logger

# 获取日志记录器
logger = get_logger(__name__)


def extract_text(file_path: str) -> str:
    """提取文档文本内容"""
    suffix = Path(file_path).suffix.lower()
    logger.info(f"开始提取文本: {file_path}")
    
    try:
        if suffix == ".pdf":
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                max_pages = min(MAX_PDF_PAGES, len(reader.pages))
                text_parts = [page.extract_text() for page in reader.pages[:max_pages] if page.extract_text()]
                result = "\n".join(text_parts).strip()
                logger.info(f"PDF 文本提取完成: {len(reader.pages)} 页，读取 {max_pages} 页，字符数: {len(result)}")
                return result
        
        elif suffix == ".docx":
            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs[:MAX_DOCX_PARAGRAPHS] if p.text.strip()]
            result = "\n".join(paragraphs)
            logger.info(f"DOCX 文本提取完成: 读取 {len(paragraphs)} 段，字符数: {len(result)}")
            return result
        
        elif suffix == ".xlsx":
            df = pd.read_excel(file_path, sheet_name=0, engine='openpyxl', nrows=MAX_XLSX_ROWS)
            
            # 格式化Excel问卷数据
            text_parts = []
            text_parts.append(f"问卷字段：{', '.join(df.columns.tolist())}")
            text_parts.append("\n" + "="*50 + "\n")
            
            # 遍历每一行（每个问卷记录）
            for idx, row in df.iterrows():
                text_parts.append(f"\n【问卷记录 {idx + 1}】")
                for col in df.columns:
                    value = row[col]
                    # 过滤空值和"无"
                    if pd.notna(value) and str(value).strip() and str(value).strip() != '无':
                        text_parts.append(f"{col}: {value}")
                text_parts.append("")
            
            result = '\n'.join(text_parts)
            logger.info(f"XLSX 文本提取完成: 读取 {len(df)} 行，{len(df.columns)} 列，字符数: {len(result)}")
            return result
        
        elif suffix == ".txt":
            with open(file_path, 'r', encoding='utf-8') as f:
                result = f.read(MAX_TXT_CHARS)
                logger.info(f"TXT 文本提取完成: 字符数: {len(result)}")
                return result
        
        elif suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            text_content = []
            max_slides = min(MAX_PPTX_SLIDES, len(prs.slides))
            for slide_idx, slide in enumerate(prs.slides):
                if slide_idx >= max_slides:
                    break
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text.strip())
            result = "\n".join(text_content)
            logger.info(f"PPTX 文本提取完成: {len(prs.slides)} 页，读取 {max_slides} 页，字符数: {len(result)}")
            return result
        
        else:
            logger.error(f"不支持的文件格式：{suffix}")
            raise ValueError(f"不支持的文件格式：{suffix}")
    
    except ImportError as e:
        if "pptx" in str(e):
            logger.error("缺少 python-pptx 库")
            raise Exception("请安装 python-pptx 库：pip install python-pptx")
        raise
    except Exception as e:
        logger.error(f"读取文件 {file_path} 失败：{str(e)}")
        raise Exception(f"读取文件 {file_path} 失败：{str(e)}")


def find_project_files(directory: str = None) -> list:
    """
    查找目录中的文档文件
    
    Args:
        directory: 目录路径，默认从 data/ 目录查找
        
    Returns:
        文件路径列表（最多返回前2个文件）
    """
    if directory is None:
        # 获取项目根目录
        from src.utils.config import BASE_DIR
        directory = BASE_DIR / "data"
    
    files = _find_files_in_dir(str(directory))
    
    if files:
        logger.info(f"在 {directory} 目录中找到 {len(files)} 个文件")
        # 只返回前2个文件
        if len(files) > 2:
            logger.info(f"系统将使用前2个文件: {files[0]}, {files[1]}")
            return files[:2]
        return files
    else:
        logger.warning(f"在 {directory} 目录中未找到任何文档文件")
        logger.info(f"提示: 请确保文件在 {directory} 目录下")
        return []


def _find_files_in_dir(directory: str) -> list:
    """
    在指定目录中查找文档文件
    
    Args:
        directory: 目录路径
        
    Returns:
        文件路径列表（按文件名排序）
    """
    path = Path(directory)
    
    # 检查目录是否存在
    if not path.exists():
        logger.warning(f"目录不存在: {directory}")
        return []
    
    files = []
    
    # 查找所有支持格式的文件
    for ext in SUPPORTED_FORMATS:
        files.extend(path.glob(f"*{ext}"))
    
    # 转换为字符串列表并按文件名排序
    files = sorted([str(f) for f in files], key=lambda x: Path(x).name)
    
    return files

